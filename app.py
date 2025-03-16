import streamlit as st
from PyPDF2 import PdfReader, PdfWriter
from PIL import Image
from io import BytesIO
from docx import Document
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import fitz  # PyMuPDF for PDF compression

st.set_page_config(page_title="PDF & File Converter", layout="wide")

# ✅ Load Custom CSS
def load_css():
    with open("assets/Style.css", "r") as css_file:
        st.markdown(f"<style>{css_file.read()}</style>", unsafe_allow_html=True)

load_css()

# ✅ Display Logo
st.image("logo1.png", width=150)
st.markdown('<p class="title">📄 PDF & File Converter</p>', unsafe_allow_html=True)

# --- Select Operation ---
operation = st.selectbox("Select an operation:", [
    "Generate Empty PDF",
    "Convert Any File to PDF",
    "Extract Pages from PDF",
    "Merge PDFs",
    "Split PDF (1 to 2 PDFs)",
    "Compress PDF",
    "Insert Page Numbers to PDF"
])

# ✅ Generate Empty PDF Functionality
if operation == "Generate Empty PDF":
    st.subheader("📄 Generate an Empty PDF")

    # ✅ Enter Number of Pages
    num_pages = st.number_input("Enter number of pages:", min_value=1, max_value=100, value=1, step=1)

    # ✅ Button to Generate PDF
    generate_btn = st.button("Generate an Empty PDF")

    if generate_btn:
        output_pdf = BytesIO()
        pdf_canvas = canvas.Canvas(output_pdf, pagesize=letter)
        pdf_canvas.setFont("Helvetica", 12)

        for i in range(num_pages):
            pdf_canvas.drawString(100, 750, f"Page {i+1}")
            pdf_canvas.showPage()

        pdf_canvas.save()
        output_pdf.seek(0)

        st.success(f"✅ Empty PDF with {num_pages} pages generated!")
        st.download_button("📥 Download Empty PDF", data=output_pdf, file_name="Empty_PDF.pdf", mime="application/pdf")

    st.stop()  # ✅ Stops execution for other operations
# ✅ File Upload Section
uploaded_files = st.file_uploader("Upload file(s)", type=["pdf", "png", "jpg", "jpeg", "docx", "pptx", "txt"], accept_multiple_files=True)

if uploaded_files:
    st.success(f"✅ {len(uploaded_files)} file(s) uploaded!")

    # ✅ Convert Any File to PDF
    if operation == "Convert Any File to PDF":
        for uploaded_file in uploaded_files:
            file_name = uploaded_file.name.rsplit(".", 1)[0]
            file_extension = uploaded_file.name.rsplit(".", 1)[-1]

            output_pdf = BytesIO()

            if file_extension in ["png", "jpg", "jpeg"]:
                Image.open(uploaded_file).convert("RGB").save(output_pdf, format="PDF")
                output_pdf.seek(0)

            elif file_extension == "txt":
                pdf_canvas = canvas.Canvas(output_pdf, pagesize=letter)
                pdf_canvas.setFont("Helvetica", 12)
                y_position = 750
                for line in uploaded_file.getvalue().decode().split("\n"):
                    pdf_canvas.drawString(100, y_position, line)
                    y_position -= 20
                pdf_canvas.save()
                output_pdf.seek(0)

            elif file_extension == "docx":
                doc = Document(uploaded_file)
                pdf_canvas = canvas.Canvas(output_pdf, pagesize=letter)
                pdf_canvas.setFont("Helvetica", 12)
                y_position = 750
                for para in doc.paragraphs:
                    pdf_canvas.drawString(100, y_position, para.text)
                    y_position -= 20
                pdf_canvas.save()
                output_pdf.seek(0)

            elif file_extension == "pptx":
                ppt = Presentation(uploaded_file)
                pdf_canvas = canvas.Canvas(output_pdf, pagesize=letter)
                pdf_canvas.setFont("Helvetica", 12)
                y_position = 750
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            pdf_canvas.drawString(100, y_position, shape.text)
                            y_position -= 20
                pdf_canvas.save()
                output_pdf.seek(0)

            else:
                st.error(f"❌ Unsupported file format: {file_extension}")
                continue

            st.download_button(f"💚 Download {file_name}.pdf", data=output_pdf, file_name=f"{file_name}.pdf", mime="application/pdf")

    # ✅ Extract Pages from PDF
    elif operation == "Extract Pages from PDF":
        pdf_reader = PdfReader(uploaded_files[0])
        pages_to_extract = st.text_input("Enter page numbers (comma-separated):")

        if st.button("Extract"):
            if pages_to_extract:
                selected_pages = [int(p.strip()) - 1 for p in pages_to_extract.split(",")]
                pdf_writer = PdfWriter()
                for p in selected_pages:
                    pdf_writer.add_page(pdf_reader.pages[p])

                output_pdf = BytesIO()
                pdf_writer.write(output_pdf)
                output_pdf.seek(0)
                st.download_button("📄 Download Extracted PDF", data=output_pdf, file_name="Extracted_Pages.pdf", mime="application/pdf")

    # ✅ Merge PDFs
    elif operation == "Merge PDFs":
        pdf_writer = PdfWriter()
        for file in uploaded_files:
            pdf_reader = PdfReader(file)
            for page in pdf_reader.pages:
                pdf_writer.add_page(page)

        output_pdf = BytesIO()
        pdf_writer.write(output_pdf)
        output_pdf.seek(0)
        st.download_button("💚 Download Merged PDF", data=output_pdf, file_name="Merged_PDF.pdf", mime="application/pdf")

    # ✅ Split PDF
    elif operation == "Split PDF (1 to 2 PDFs)":
        pdf_reader = PdfReader(uploaded_files[0])
        total_pages = len(pdf_reader.pages)

        if total_pages > 1:
            mid = total_pages // 2
            part1_writer, part2_writer = PdfWriter(), PdfWriter()

            for i in range(mid):
                part1_writer.add_page(pdf_reader.pages[i])
            for i in range(mid, total_pages):
                part2_writer.add_page(pdf_reader.pages[i])

            output1, output2 = BytesIO(), BytesIO()
            part1_writer.write(output1)
            part2_writer.write(output2)
            output1.seek(0)
            output2.seek(0)

            st.download_button("📄 Download First Half", data=output1, file_name="Part1.pdf", mime="application/pdf")
            st.download_button("📄 Download Second Half", data=output2, file_name="Part2.pdf", mime="application/pdf")
        else:
            st.error("❌ The PDF must have at least 2 pages to split.")

    # ✅ Compress PDF
    elif operation == "Compress PDF":
        pdf_reader = fitz.open(stream=uploaded_files[0].read(), filetype="pdf")
        output_pdf = BytesIO()

        for page in pdf_reader:
            page.clean_contents()

        pdf_reader.save(output_pdf)
        output_pdf.seek(0)
        st.download_button("📉 Download Compressed PDF", data=output_pdf, file_name="Compressed_PDF.pdf", mime="application/pdf")

    # ✅ Insert Page Numbers
    elif operation == "Insert Page Numbers to PDF":
        pdf_reader = PdfReader(uploaded_files[0])
        pdf_writer = PdfWriter()
        output_pdf = BytesIO()

        for i, page in enumerate(pdf_reader.pages):
            packet = BytesIO()
            c = canvas.Canvas(packet, pagesize=letter)
            c.setFont("Helvetica", 12)
            c.drawString(500, 20, f"Page {i + 1}")
            c.save()

            packet.seek(0)
            overlay_reader = PdfReader(packet)
            page.merge_page(overlay_reader.pages[0])
            pdf_writer.add_page(page)

        pdf_writer.write(output_pdf)
        output_pdf.seek(0)
        st.download_button("📄 Download Numbered PDF", data=output_pdf, file_name="Numbered_PDF.pdf", mime="application/pdf")

# ✅ Copyright
st.markdown('<p class="small-text">© Pavan Sri Sai Mondem | Siva Satyamsetti | Uma Satyam Mounika Sapireddy | Bhuvaneswari Devi Seru | Chandu Meela | Trainees from Techwing 🧡</p>', unsafe_allow_html=True)
